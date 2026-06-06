# SPDX-License-Identifier: MIT
"""
测试报告生成 - Word + 通知（企业微信/飞书/钉钉）
被引用方：09-报告生成 agent
"""
import json
import logging
import os
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional

import requests

sys.path.insert(0, str(Path(__file__).parent.parent))
from utils.paths import get_output_dir, current_run_id

logger = logging.getLogger(__name__)


# ===== Word 报告 =====

def _write_docx_header(doc, data: Dict) -> None:
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    style = doc.styles["Normal"]
    style.font.name = "微软雅黑"
    style.font.size = Pt(11)
    title = doc.add_heading("测试报告", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"项目：{data.get('project_name', '')}")
    doc.add_paragraph(f"版本：{data.get('version', '')}")
    doc.add_paragraph(f"测试环境：{data.get('environment', '')}")
    doc.add_paragraph(f"报告日期：{datetime.now().strftime('%Y-%m-%d %H:%M')}")
    doc.add_page_break()


def _write_docx_summary(doc, data: Dict) -> None:
    from docx.shared import RGBColor
    doc.add_heading("一、执行摘要", level=1)
    summary_table = doc.add_table(rows=2, cols=5)
    summary_table.style = "Table Grid"
    for i, h in enumerate(["总用例", "通过", "失败", "通过率", "覆盖率"]):
        cell = summary_table.cell(0, i)
        cell.text = h
        if cell.paragraphs and cell.paragraphs[0].runs:
            cell.paragraphs[0].runs[0].bold = True
    results = data.get("results", {})
    for i, v in enumerate([
        str(results.get("total", 0)), str(results.get("passed", 0)),
        str(results.get("failed", 0)), f"{results.get('pass_rate', 0):.1%}",
        f"{data.get('coverage', 0):.1%}",
    ]):
        summary_table.cell(1, i).text = v
    conclusion = doc.add_paragraph()
    conclusion.add_run("测试结论：").bold = True
    verdict = data.get("verdict", "通过")
    run = conclusion.add_run(verdict)
    run.font.color.rgb = RGBColor(0, 128, 0) if verdict == "通过" else RGBColor(255, 0, 0)
    run.bold = True


def _write_docx_degraded_warning(doc, data: Dict) -> None:
    from docx.shared import RGBColor
    degraded_upstream = data.get("_degraded_upstream", [])
    if not degraded_upstream:
        return
    doc.add_heading("⚠ 数据完整性警示", level=1)
    warning_p = doc.add_paragraph()
    warning_run = warning_p.add_run(
        f"本次报告基于不完整测试数据生成。共 {len(degraded_upstream)} 个 expert "
        f"输出 degraded(mock 兜底 / LLM 失败 / 未实装 V1.x rollout):"
    )
    warning_run.font.color.rgb = RGBColor(255, 140, 0)
    warning_run.bold = True
    for name in degraded_upstream:
        item = doc.add_paragraph(style="List Bullet")
        item.add_run(f"expert '{name}' — 详见 ROADMAP.md V1.x rollout 节奏")
    impact_p = doc.add_paragraph()
    impact_run = impact_p.add_run("→ 上线决策建议: conditional 或 no-go(由 test-lead 判定);不应基于此报告直接发版。")
    impact_run.bold = True
    impact_run.font.color.rgb = RGBColor(255, 0, 0)


def _write_docx_bugs(doc, data: Dict) -> None:
    doc.add_heading("二、缺陷统计", level=1)
    bugs = data.get("bugs", {})
    bug_table = doc.add_table(rows=5, cols=3)
    bug_table.style = "Table Grid"
    rows = [
        ("级别", "数量", "修复率"),
        ("P0", str(bugs.get("p0", 0)), f"{bugs.get('p0_fix_rate', 0):.0%}"),
        ("P1", str(bugs.get("p1", 0)), f"{bugs.get('p1_fix_rate', 0):.0%}"),
        ("P2", str(bugs.get("p2", 0)), f"{bugs.get('p2_fix_rate', 0):.0%}"),
        ("P3", str(bugs.get("p3", 0)), f"{bugs.get('p3_fix_rate', 0):.0%}"),
    ]
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = bug_table.cell(ri, ci)
            cell.text = val
            if ri == 0 and cell.paragraphs and cell.paragraphs[0].runs:
                cell.paragraphs[0].runs[0].bold = True


def _write_docx_performance(doc, data: Dict) -> None:
    perf = data.get("performance")
    if not perf:
        return
    doc.add_heading("三、性能指标（JMeter）", level=1)
    p_table = doc.add_table(rows=5, cols=2)
    p_table.style = "Table Grid"
    for ri, (k, v) in enumerate([
        ("TPS", f"{perf.get('tps', 0)} 次/秒"),
        ("平均响应", f"{perf.get('avg_response_ms', 0)} ms"),
        ("P95响应", f"{perf.get('p95_response_ms', 0)} ms"),
        ("错误率", f"{perf.get('error_rate_pct', 0)}%"),
        ("门禁结论", perf.get("quality_gate", "PASS")),
    ]):
        p_table.cell(ri, 0).text = k
        p_table.cell(ri, 1).text = v


def _write_docx_risks(doc, data: Dict) -> None:
    doc.add_heading("四、风险与建议", level=1)
    for risk in data.get("risks", []):
        para = doc.add_paragraph(style="List Bullet")
        para.add_run(f"【{risk.get('level', '中')}】").bold = True
        para.add_run(risk.get("description", ""))


def generate_test_report(data: Dict, output_path: str) -> str:
    """生成 Word 测试报告。

    依赖 python-docx (可选)。未装时 graceful skip + sentinel 文件标记。
    """
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx 未安装, Word 报告跳过生成。pip install python-docx>=1.1.0")
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        sentinel = Path(output_path).with_suffix(".skipped.txt")
        results = data.get("results", {})
        sentinel.write_text(
            f"Word report skipped: python-docx not installed.\n"
            f"project={data.get('project_name', '-')} verdict={data.get('verdict', '-')} "
            f"pass_rate={results.get('pass_rate', 0):.1%}\n"
            "Install python-docx>=1.1.0.\n",
            encoding="utf-8",
        )
        logger.info(f"sentinel 文件: {sentinel}")
        return str(sentinel)

    doc = Document()
    _write_docx_header(doc, data)
    _write_docx_summary(doc, data)
    _write_docx_degraded_warning(doc, data)
    _write_docx_bugs(doc, data)
    _write_docx_performance(doc, data)
    _write_docx_risks(doc, data)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    doc.save(output_path)
    logger.info(f"Word 报告已生成: {output_path}")
    return output_path


# ===== 通知 =====

def send_wechat_report(summary: Dict, webhook: Optional[str] = None) -> bool:
    """企业微信 markdown 通知"""
    webhook = webhook or os.getenv("WECHAT_WEBHOOK_URL")
    if not webhook:
        logger.warning("未配置 WECHAT_WEBHOOK_URL，跳过企业微信通知")
        return False
    pass_rate = summary.get("pass_rate", 0)
    verdict = summary.get("verdict", "通过")
    icon = "✅" if verdict == "通过" else "❌"
    content = (
        f"{icon} **{summary.get('project','')} 测试完成**\n"
        f"> 环境：{summary.get('environment','')}\n"
        f"> 用例：{summary.get('total',0)} 通过 {summary.get('passed',0)}（{pass_rate:.1%}）\n"
        f"> Bug：P0({summary.get('p0_bugs',0)}) P1({summary.get('p1_bugs',0)})\n"
        f"> 性能：TPS={summary.get('perf_tps','-')} P95={summary.get('perf_p95','-')}ms\n"
        f"> [查看完整报告]({summary.get('report_url','#')})"
    )
    try:
        r = requests.post(webhook, json={"msgtype": "markdown", "markdown": {"content": content}}, timeout=10)
        return r.ok
    except Exception as e:
        logger.error(f"企业微信通知失败: {e}")
        return False


def send_feishu_report(summary: Dict, webhook: Optional[str] = None) -> bool:
    """飞书 富文本卡片 通知"""
    webhook = webhook or os.getenv("FEISHU_WEBHOOK_URL")
    if not webhook:
        logger.warning("未配置 FEISHU_WEBHOOK_URL，跳过飞书通知")
        return False
    # 飞书卡片合法颜色：blue/wathet/turquoise/green/yellow/orange/red/carmine/violet/purple/indigo/grey
    color = "green" if summary.get("pass_rate", 0) >= 0.9 else "red"
    card = {
        "config": {"wide_screen_mode": True},
        "header": {
            "title": {"content": f"测试报告 - {summary.get('project','')}", "tag": "plain_text"},
            "template": color,
        },
        "elements": [
            {
                "tag": "div",
                "fields": [
                    {"is_short": True, "text": {"tag": "lark_md", "content": f"**环境**\n{summary.get('environment','')}"}},
                    {"is_short": True, "text": {"tag": "lark_md", "content": f"**通过率**\n{summary.get('pass_rate',0):.1%}"}},
                    {"is_short": True, "text": {"tag": "lark_md", "content": f"**P0 Bug**\n{summary.get('p0_bugs',0)}个"}},
                    {"is_short": True, "text": {"tag": "lark_md", "content": f"**覆盖率**\n{summary.get('coverage',0):.1%}"}},
                ],
            },
            {
                "tag": "action",
                "actions": [{
                    "tag": "button",
                    "text": {"content": "查看完整报告", "tag": "plain_text"},
                    "url": summary.get("report_url", "#"),
                    "type": "primary",
                }],
            },
        ],
    }
    try:
        r = requests.post(webhook, json={"msg_type": "interactive", "card": card}, timeout=10)
        return r.ok
    except Exception as e:
        logger.error(f"飞书通知失败: {e}")
        return False


def send_dingtalk_report(summary: Dict, webhook: Optional[str] = None) -> bool:
    """钉钉 markdown 通知"""
    webhook = webhook or os.getenv("DINGTALK_WEBHOOK_URL")
    if not webhook:
        logger.warning("未配置 DINGTALK_WEBHOOK_URL，跳过钉钉通知")
        return False
    pass_rate = summary.get("pass_rate", 0)
    verdict = summary.get("verdict", "通过")
    icon = "✅" if verdict == "通过" else "❌"
    text = (
        f"### {icon} {summary.get('project','')} 测试完成\n"
        f"- 环境：{summary.get('environment','')}\n"
        f"- 通过率：{pass_rate:.1%}\n"
        f"- Bug：P0({summary.get('p0_bugs',0)}) P1({summary.get('p1_bugs',0)})\n"
        f"- 性能：TPS={summary.get('perf_tps','-')} P95={summary.get('perf_p95','-')}ms\n"
        f"- [查看报告]({summary.get('report_url','#')})"
    )
    try:
        r = requests.post(webhook, json={
            "msgtype": "markdown",
            "markdown": {"title": "测试报告", "text": text},
        }, timeout=10)
        return r.ok
    except Exception as e:
        logger.error(f"钉钉通知失败: {e}")
        return False


def send_all_notifications(summary: Dict):
    """同时发送企业微信/飞书/钉钉/Server酱（按 .env 配置自动跳过未配置项）"""
    return {
        "wechat": send_wechat_report(summary),
        "feishu": send_feishu_report(summary),
        "dingtalk": send_dingtalk_report(summary),
        "serverchan": send_serverchan_report(summary),
    }


def send_serverchan_report(summary: Dict, sendkey: Optional[str] = None) -> bool:
    """Server酱(https://sct.ftqq.com/) 微信推送通知。
    SendKey 从 .env SERVERCHAN_SENDKEY 读取。
    """
    sendkey = sendkey or os.getenv("SERVERCHAN_SENDKEY", "")
    if not sendkey:
        logger.warning("未配置 SERVERCHAN_SENDKEY，跳过 Server酱通知")
        return False
    try:
        title = f"{summary.get('verdict','?')} {summary.get('project','')} {summary.get('pass_rate',0)*100:.0f}%"
        desp = (
            f"## {summary.get('project','')} 测试报告\n\n"
            f"> 结论：**{summary.get('verdict','?')}**\n"
            f"> 用例：{summary.get('total',0)} 通过 {summary.get('passed',0)}（{summary.get('pass_rate',0)*100:.1f}%）\n"
            f"> Bug：P0({summary.get('p0_bugs',0)}) P1({summary.get('p1_bugs',0)})\n"
            f"> 覆盖率：{summary.get('coverage',0)*100:.0f}%\n"
        )
        if summary.get('perf_tps'):
            desp += f"> 性能：TPS={summary['perf_tps']} P95={summary.get('perf_p95','?')}ms\n"
        if summary.get('report_url'):
            desp += f"\n[查看完整报告]({summary['report_url']})"
        r = requests.post(
            f"https://sctapi.ftqq.com/{sendkey}.send",
            data={"title": title, "desp": desp},
            timeout=10,
        )
        ok = r.json().get("code") == 0
        logger.info(f"Server酱通知: {'成功' if ok else '失败'}")
        return ok
    except Exception as e:
        logger.error(f"Server酱通知失败: {e}")
        return False


# ===== PDF 报告（reportlab） =====

def generate_pdf_report(data: Dict, output_path: str) -> str:
    """生成 PDF 测试报告。需 pip install reportlab"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    except ImportError:
        raise RuntimeError("reportlab 未安装：pip install reportlab")

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(output_path, pagesize=A4,
                            topMargin=2 * cm, bottomMargin=2 * cm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("title", parent=styles["Title"], alignment=1)

    story = []
    story.append(Paragraph(f"测试报告 - {data.get('project_name', '')}", title_style))
    story.append(Spacer(1, 0.5 * cm))
    story.append(Paragraph(f"版本: {data.get('version', '')}", styles["Normal"]))
    story.append(Paragraph(f"环境: {data.get('environment', '')}", styles["Normal"]))
    story.append(Paragraph(f"日期: {datetime.now():%Y-%m-%d %H:%M}", styles["Normal"]))
    story.append(Spacer(1, 1 * cm))

    # 摘要表
    results = data.get("results", {})
    summary_data = [
        ["指标", "数值"],
        ["总用例", str(results.get("total", 0))],
        ["通过", str(results.get("passed", 0))],
        ["失败", str(results.get("failed", 0))],
        ["通过率", f"{results.get('pass_rate', 0):.1%}"],
        ["覆盖率", f"{data.get('coverage', 0):.1%}"],
        ["结论", data.get("verdict", "通过")],
    ]
    table = Table(summary_data, colWidths=[5 * cm, 8 * cm])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#366092")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("ALIGN", (0, 0), (-1, -1), "LEFT"),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
    ]))
    story.append(table)
    story.append(Spacer(1, 0.8 * cm))

    # 性能段
    perf = data.get("performance")
    if perf:
        story.append(Paragraph("性能指标（JMeter）", styles["Heading2"]))
        perf_data = [
            ["指标", "实测", "门禁"],
            ["TPS", str(perf.get("tps", "-")), "≥100"],
            ["平均响应", f"{perf.get('avg_response_ms', '-')} ms", "≤200ms"],
            ["P95 响应", f"{perf.get('p95_response_ms', '-')} ms", "≤500ms"],
            ["错误率", f"{perf.get('error_rate_pct', '-')}%", "<1%"],
        ]
        story.append(Table(perf_data, colWidths=[4 * cm, 4 * cm, 4 * cm]))
        story.append(Spacer(1, 0.5 * cm))

    # 风险
    risks = data.get("risks", [])
    if risks:
        story.append(Paragraph("风险与建议", styles["Heading2"]))
        for r in risks:
            story.append(Paragraph(f"• 【{r.get('level', '中')}】 {r.get('description', '')}", styles["Normal"]))

    doc.build(story)
    logger.info(f"PDF 报告已生成: {output_path}")
    return output_path


# ===== PPTX 摘要（高管汇报）=====

def generate_pptx_summary(data: Dict, output_path: str) -> str:
    """生成 PPT 摘要（5 页：封面/摘要/缺陷/性能/风险与建议）"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise RuntimeError("python-pptx 未安装：pip install python-pptx")

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{data.get('project_name', '')} 测试报告"
    slide.placeholders[1].text = (
        f"版本: {data.get('version', '')}\n"
        f"环境: {data.get('environment', '')}\n"
        f"日期: {datetime.now():%Y-%m-%d}"
    )

    # 执行摘要
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "执行摘要"
    results = data.get("results", {})
    text_frame = slide.shapes.placeholders[0].text_frame
    text_frame.text = (
        f"总用例: {results.get('total', 0)}    "
        f"通过率: {results.get('pass_rate', 0):.1%}    "
        f"覆盖率: {data.get('coverage', 0):.1%}\n"
        f"结论: {data.get('verdict', '通过')}"
    )

    # 缺陷
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "缺陷统计"
    bugs = data.get("bugs", {})
    slide.shapes.placeholders[0].text_frame.text = (
        f"P0: {bugs.get('p0', 0)} (修复率 {bugs.get('p0_fix_rate', 0):.0%})\n"
        f"P1: {bugs.get('p1', 0)} (修复率 {bugs.get('p1_fix_rate', 0):.0%})\n"
        f"P2: {bugs.get('p2', 0)}    P3: {bugs.get('p3', 0)}"
    )

    # 性能
    perf = data.get("performance")
    if perf:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "性能指标"
        slide.shapes.placeholders[0].text_frame.text = (
            f"TPS: {perf.get('tps', '-')}\n"
            f"P95: {perf.get('p95_response_ms', '-')} ms\n"
            f"错误率: {perf.get('error_rate_pct', '-')}%\n"
            f"门禁: {perf.get('quality_gate', 'PASS')}"
        )

    # 风险
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "风险与建议"
    risks_text = "\n".join(
        f"• 【{r.get('level', '中')}】 {r.get('description', '')}"
        for r in data.get("risks", [])
    ) or "无重大风险"
    slide.shapes.placeholders[0].text_frame.text = risks_text

    prs.save(output_path)
    logger.info(f"PPTX 摘要已生成: {output_path}")
    return output_path


# ===== CLI =====

def main():
    import argparse
    logging.basicConfig(level=logging.INFO)
    parser = argparse.ArgumentParser(description="生成 Word 测试报告 + 通知")
    parser.add_argument("--data", required=True, help="测试结果 JSON 文件路径")
    parser.add_argument("--output", default=None, help="Word 报告输出路径")
    parser.add_argument("--notify", action="store_true", help="发送通知")
    args = parser.parse_args()

    with open(args.data, encoding="utf-8") as f:
        data = json.load(f)

    output = args.output or str(get_output_dir("reports", current_run_id()) / f"测试报告_{datetime.now().strftime('%Y%m%d')}.docx")
    generate_test_report(data, output)

    if args.notify:
        results = send_all_notifications(data)
        print(f"通知发送结果：{results}")


if __name__ == "__main__":
    main()
