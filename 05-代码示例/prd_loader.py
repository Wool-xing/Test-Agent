"""
PRD 多格式加载器 - 统一入口，支持 md/txt/pdf/word/excel/zip/img/url
被引用方：02-需求分析 agent / test-coordinator skill / test-lead 自动路由
"""
import io
import json
import logging
import os
import re
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Dict, List, Optional, Union
from urllib.parse import urlparse

logger = logging.getLogger(__name__)


# ===== 通用加载入口 =====

def load_prd(source: Union[str, Path]) -> Dict:
    """
    统一 PRD 加载入口。返回：
      {
        "format": "md|txt|pdf|docx|xlsx|zip|image|url|html",
        "source": "原始路径或 URL",
        "text": "提取的纯文本",
        "attachments": [...]   # zip 解包后的子文件列表
        "images": [...]        # PRD 含图时图像路径列表
        "metadata": {...}      # 标题/作者/页数等
      }
    """
    s = str(source)

    # URL（confluence / wiki / 链接）
    if re.match(r"^https?://", s):
        return _load_url(s)

    p = Path(s)
    if not p.exists():
        raise FileNotFoundError(f"PRD 文件不存在: {s}")

    suffix = p.suffix.lower()
    if suffix in (".md", ".markdown", ".txt"):
        return _load_text(p)
    if suffix == ".pdf":
        return _load_pdf(p)
    if suffix in (".doc", ".docx"):
        return _load_docx(p)
    if suffix in (".xls", ".xlsx"):
        return _load_xlsx(p)
    if suffix == ".zip":
        return _load_zip(p)
    if suffix in (".png", ".jpg", ".jpeg", ".bmp", ".webp"):
        return _load_image(p)
    if suffix in (".html", ".htm"):
        return _load_html(p)
    if suffix == ".pptx":
        return _load_pptx(p)

    # 兜底当文本读
    logger.warning(f"未知格式 {suffix}，按文本读取")
    return _load_text(p)


# ===== md / txt =====

def _load_text(p: Path) -> Dict:
    text = p.read_text(encoding="utf-8", errors="replace")
    return {
        "format": p.suffix.lstrip(".").lower(),
        "source": str(p),
        "text": text,
        "attachments": [],
        "images": [],
        "metadata": {"size_bytes": p.stat().st_size, "name": p.name},
    }


# ===== PDF =====

def _load_pdf(p: Path) -> Dict:
    """PDF 解析：优先 pdfplumber（保留布局），fallback pypdf"""
    text = ""
    metadata = {}
    images: List[str] = []

    try:
        import pdfplumber
        with pdfplumber.open(p) as pdf:
            metadata["pages"] = len(pdf.pages)
            for page in pdf.pages:
                t = page.extract_text() or ""
                text += t + "\n"
            # PDF 中嵌入的图片暂不抽取（需 pdf2image + poppler，按需启用）
        logger.info(f"PDF 解析（pdfplumber）: {p.name}, {metadata['pages']} 页")
    except ImportError:
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(p))
            metadata["pages"] = len(reader.pages)
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
            logger.info(f"PDF 解析（pypdf fallback）: {p.name}")
        except ImportError:
            raise RuntimeError("PDF 解析需要 pdfplumber 或 pypdf: pip install pdfplumber")

    return {
        "format": "pdf", "source": str(p), "text": text.strip(),
        "attachments": [], "images": images,
        "metadata": metadata,
    }


# ===== Word docx =====

def _load_docx(p: Path) -> Dict:
    """docx 解析：python-docx（已在 requirements 中）"""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("docx 解析需要 python-docx: pip install python-docx")

    doc = Document(str(p))
    paras = [para.text for para in doc.paragraphs if para.text.strip()]
    text = "\n".join(paras)

    # 表格内容也提取
    tables_text = []
    for tbl in doc.tables:
        for row in tbl.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            tables_text.append(row_text)
    if tables_text:
        text += "\n\n=== 表格 ===\n" + "\n".join(tables_text)

    metadata = {
        "title": doc.core_properties.title or "",
        "author": doc.core_properties.author or "",
        "created": str(doc.core_properties.created) if doc.core_properties.created else "",
        "paragraphs": len(paras),
        "tables": len(doc.tables),
    }
    logger.info(f"docx 解析: {p.name}, {len(paras)} 段, {len(doc.tables)} 表")
    return {
        "format": "docx", "source": str(p), "text": text,
        "attachments": [], "images": [], "metadata": metadata,
    }


# ===== Excel xlsx =====

def _load_xlsx(p: Path) -> Dict:
    """Excel 解析：每个 Sheet 转 markdown 表格"""
    import openpyxl
    wb = openpyxl.load_workbook(p, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_parts.append(f"## Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join("" if v is None else str(v) for v in row)
            if row_text.strip(" |"):
                text_parts.append(row_text)
        text_parts.append("")
    text = "\n".join(text_parts)
    metadata = {"sheets": wb.sheetnames}
    logger.info(f"xlsx 解析: {p.name}, sheets={wb.sheetnames}")
    return {
        "format": "xlsx", "source": str(p), "text": text,
        "attachments": [], "images": [], "metadata": metadata,
    }


# ===== ZIP =====

def _load_zip(p: Path) -> Dict:
    """zip 解包到临时目录，递归加载每个子文件，合并文本"""
    extract_dir = Path(tempfile.mkdtemp(prefix="prd_zip_"))
    with zipfile.ZipFile(p, "r") as zf:
        zf.extractall(extract_dir)

    attachments: List[Dict] = []
    images: List[str] = []
    text_parts: List[str] = []

    for child in sorted(extract_dir.rglob("*")):
        if child.is_file():
            try:
                child_info = load_prd(child)
                attachments.append({
                    "name": str(child.relative_to(extract_dir)),
                    "format": child_info["format"],
                    "size_bytes": child.stat().st_size,
                })
                if child_info["format"] == "image":
                    images.append(str(child))
                elif child_info.get("text"):
                    text_parts.append(f"## 子文件: {child.relative_to(extract_dir)}\n{child_info['text']}")
            except Exception as e:
                logger.warning(f"zip 子文件加载失败 {child}: {e}")

    text = "\n\n".join(text_parts)
    logger.info(f"zip 解析: {p.name}, 含 {len(attachments)} 个子文件")
    return {
        "format": "zip", "source": str(p), "text": text,
        "attachments": attachments, "images": images,
        "metadata": {"extract_dir": str(extract_dir), "file_count": len(attachments)},
    }


# ===== 图片（需视觉模型读，本工具仅返回路径，由 Claude Code 直接看图） =====

def _load_image(p: Path) -> Dict:
    """图片 PRD：返回路径，提示 Claude Code 用视觉能力读取"""
    return {
        "format": "image", "source": str(p),
        "text": f"[图片 PRD：{p.name}，需 Claude Code 视觉能力解读，或先用 utils.visual_helper.ocr_image 转文字]",
        "attachments": [], "images": [str(p)],
        "metadata": {"size_bytes": p.stat().st_size, "name": p.name},
    }


# ===== PPTX =====

def _load_pptx(p: Path) -> Dict:
    """python-pptx 解析 PPT 幻灯片文本"""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx 未安装：pip install python-pptx")

    prs = Presentation(str(p))
    text_parts = []
    slides_count = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        text_parts.append(f"## 幻灯片 {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)
        text_parts.append("")

    return {
        "format": "pptx", "source": str(p), "text": "\n".join(text_parts),
        "attachments": [], "images": [],
        "metadata": {"slides": slides_count, "name": p.name},
    }


# ===== HTML =====

def _load_html(p: Path) -> Dict:
    text = p.read_text(encoding="utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(text, "html.parser")
        plain = soup.get_text(separator="\n", strip=True)
    except ImportError:
        plain = re.sub(r"<[^>]+>", "", text)
    return {
        "format": "html", "source": str(p), "text": plain,
        "attachments": [], "images": [], "metadata": {"name": p.name},
    }


# ===== URL（confluence / wiki / 在线文档） =====

def _load_url(url: str) -> Dict:
    """
    URL 加载：用 requests 抓 HTML 后提取正文。
    注：confluence/notion 等私有空间需先认证，建议手动下载为 PDF/HTML 再加载。
    """
    import requests
    headers = {"User-Agent": "Mozilla/5.0 PRDLoader/1.0"}
    if os.getenv("PRD_HTTP_TOKEN"):
        headers["Authorization"] = f"Bearer {os.environ['PRD_HTTP_TOKEN']}"

    r = requests.get(url, headers=headers, timeout=20)
    r.raise_for_status()

    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(r.text, "html.parser")
        # 删 script/style/nav
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        title = soup.title.string if soup.title else ""
    except ImportError:
        text = re.sub(r"<[^>]+>", "", r.text)
        title = ""

    return {
        "format": "url", "source": url, "text": text,
        "attachments": [], "images": [],
        "metadata": {"url": url, "title": title, "host": urlparse(url).hostname},
    }


# ===== 产品形态识别（test-lead 路由用） =====

PLATFORM_KEYWORDS = {
    "mobile_android": [
        "android", "安卓", ".apk", "appium", "google play", "kotlin", "adb",
    ],
    "mobile_ios": [
        "ios", "ipa", "iphone", "ipad", "xcode", "swift", "testflight", "app store",
    ],
    "miniprogram_wx": [
        "微信小程序", "wechat miniprogram", "mp.weixin", "wxss", "wxml", "wx.app",
    ],
    "miniprogram_alipay": [
        "支付宝小程序", "alipay miniprogram",
    ],
    "desktop_windows": [
        ".exe", "windows 桌面", "winforms", "wpf", "win32", "uwp",
        "客户端 windows", "msi 安装包",
    ],
    "desktop_macos": [
        ".app", "macos", "mac 客户端", "cocoa", "swift mac", "dmg",
    ],
    "desktop_linux": [
        "linux 客户端", "linux 桌面", "gtk", "deb", "rpm",
    ],
    "electron": [
        "electron", "vscode", "钉钉 pc", "飞书 pc", "im 客户端 pc",
    ],
    "web": [
        "网页", "web", "h5", "网站", "spa", "react", "vue", "angular", "前端",
    ],
    "api": [
        "接口", "api", "rest", "graphql", "微服务", "后端 api",
    ],
    "performance": [
        "性能", "压测", "tps", "qps", "并发", "高并发", "吞吐量", "延迟",
    ],
    "game_visual": [
        "游戏", "canvas", "webgl", "unity", "unreal", "ue", "cocos", "图像识别",
    ],
    "iot": [
        "iot", "嵌入式", "串口", "mqtt", "modbus", "硬件设备", "传感器",
    ],
    "media": [
        "音视频", "视频", "音频", "直播", "rtsp", "rtmp", "hls", "ffmpeg",
    ],
    "tracing_mq": [
        "链路追踪", "jaeger", "zipkin", "kafka", "rabbitmq", "rocketmq", "消息队列",
    ],
    "ai_ml": [
        "模型", "ai", "ml", "llm", "chatgpt", "gpt", "claude", "推理", "推荐算法",
        "分类模型", "回归模型", "深度学习",
    ],
}

# 平台 → 推荐 agent + skill
PLATFORM_ROUTING = {
    "mobile_android": {"agent": "mobile-tester", "skill": "/mobile-test"},
    "mobile_ios": {"agent": "mobile-tester", "skill": "/mobile-test"},
    "miniprogram_wx": {"agent": "mobile-tester", "skill": "/mobile-test"},
    "miniprogram_alipay": {"agent": "mobile-tester", "skill": "/mobile-test"},
    "desktop_windows": {"agent": "desktop-tester", "skill": "/desktop-test"},
    "desktop_macos": {"agent": "desktop-tester", "skill": "/desktop-test"},
    "desktop_linux": {"agent": "desktop-tester", "skill": "/desktop-test"},
    "electron": {"agent": "desktop-tester", "skill": "/desktop-test"},
    "web": {"agent": "automation-engineer", "skill": "/python-script-gen"},
    "api": {"agent": "automation-engineer", "skill": "/python-script-gen"},
    "performance": {"agent": "automation-engineer", "skill": "/jmeter-script-gen"},
    "game_visual": {"agent": "visual-tester", "skill": "/visual-test"},
    "iot": {"agent": "system-tester", "skill": "/system-test"},
    "media": {"agent": "system-tester", "skill": "/system-test"},
    "tracing_mq": {"agent": "system-tester", "skill": "/system-test"},
    "ai_ml": {"agent": "ai-tester", "skill": "/ai-test"},
}


def detect_platforms(text: str) -> List[str]:
    """从 PRD 文本中识别涉及的产品形态（多形态可并存）"""
    lower = text.lower()
    detected = []
    for platform, keywords in PLATFORM_KEYWORDS.items():
        if any(k.lower() in lower for k in keywords):
            detected.append(platform)
    return detected


def suggest_agents(text: str) -> Dict:
    """
    返回路由建议：
      {"platforms": [...], "agents": [...], "skills": [...]}
    """
    platforms = detect_platforms(text)
    agents = sorted({PLATFORM_ROUTING[p]["agent"] for p in platforms if p in PLATFORM_ROUTING})
    skills = sorted({PLATFORM_ROUTING[p]["skill"] for p in platforms if p in PLATFORM_ROUTING})
    return {
        "platforms": platforms,
        "recommended_agents": agents,
        "recommended_skills": skills,
    }


# ===== CLI =====

def main():
    import argparse
    logging.basicConfig(level=logging.INFO)
    parser = argparse.ArgumentParser(description="PRD 多格式加载 + 平台识别")
    parser.add_argument("source", help="文件路径 或 URL")
    parser.add_argument("--detect", action="store_true", help="同时输出平台识别结果")
    parser.add_argument("--save-text", help="保存提取文本到指定路径")
    args = parser.parse_args()

    info = load_prd(args.source)
    print(f"格式: {info['format']}")
    print(f"来源: {info['source']}")
    print(f"文本长度: {len(info['text'])} chars")
    print(f"附件: {len(info['attachments'])}, 图片: {len(info['images'])}")
    print(f"metadata: {info['metadata']}")

    if args.save_text:
        Path(args.save_text).parent.mkdir(parents=True, exist_ok=True)
        Path(args.save_text).write_text(info["text"], encoding="utf-8")
        print(f"文本已保存: {args.save_text}")

    if args.detect:
        routing = suggest_agents(info["text"])
        print("\n=== 平台识别 ===")
        print(json.dumps(routing, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
